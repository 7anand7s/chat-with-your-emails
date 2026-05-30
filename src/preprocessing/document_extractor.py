"""Extract content from document attachments.

Strategy:
- Text files → direct read
- PDFs → text extraction + convert pages to images for VLM
- DOCX → text extraction + convert to images for VLM
- XLSX/PPTX → text extraction
- Images → pass through to VLM
- Encrypted PDFs → decrypt with password from env
"""

import csv
import io
import json
import os
import zipfile
from dataclasses import dataclass

from config.settings import config


@dataclass
class ExtractedDocument:
    """Result of document extraction."""
    text: str  # Extracted text content
    images: list[bytes]  # Page images as PNG bytes (for VLM)
    metadata: dict  # File metadata
    is_scanned: bool = False  # True if likely image-based (needs VLM)
    page_count: int = 0


class DocumentExtractor:
    """Extract content from all document types."""

    TEXT_TYPES = {
        "text/plain", "text/csv", "text/html", "text/xml", "text/markdown",
        "text/x-python", "text/x-java", "text/x-c", "text/x-script",
        "application/json", "application/xml", "application/javascript",
        "application/x-yaml", "application/yaml", "text/yaml",
    }

    TEXT_EXTENSIONS = {
        ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json",
        ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
        ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
        ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh",
        ".html", ".htm", ".css", ".scss", ".less",
        ".sql", ".r", ".scala", ".kt", ".swift",
        ".log", ".out", ".err", ".env", ".gitignore",
        ".dockerfile", ".makefile", ".cmake",
        ".vcf", ".ics",
    }

    # PDF passwords from environment
    PDF_PASSWORD_ENV = "PDF_PASSWORDS"  # Comma-separated passwords

    def __init__(self):
        self.pdf_passwords = self._load_pdf_passwords()

    def _load_pdf_passwords(self) -> list[str]:
        """Load PDF passwords from environment variable."""
        passwords_env = os.environ.get(self.PDF_PASSWORD_ENV, "")
        if passwords_env:
            return [p.strip() for p in passwords_env.split(",") if p.strip()]
        return []

    def extract(self, filename: str, mime_type: str, data: bytes) -> ExtractedDocument | None:
        """Extract content from an attachment.

        Returns ExtractedDocument or None if type is not supported
        (images should go to VLM directly).
        """
        if len(data) > config.max_attachment_size:
            return ExtractedDocument(
                text=f"[File too large: {len(data):,} bytes]",
                images=[],
                metadata={"filename": filename, "size": len(data), "error": "too_large"},
            )

        ext = os.path.splitext(filename.lower())[1]

        # Text files
        if mime_type in self.TEXT_TYPES or ext in self.TEXT_EXTENSIONS:
            return self._extract_text_file(data, filename)

        # PDF
        if mime_type == "application/pdf" or ext == ".pdf":
            return self._extract_pdf(data, filename)

        # DOCX
        if ext in (".docx", ".doc"):
            return self._extract_docx(data, filename)

        # XLSX
        if ext in (".xlsx", ".xls"):
            return self._extract_xlsx(data, filename)

        # PPTX
        if ext in (".pptx", ".ppt"):
            return self._extract_pptx(data, filename)

        # ZIP
        if ext in (".zip", ".jar", ".war", ".apk"):
            return self._extract_zip(data, filename)

        # EML
        if ext in (".eml", ".msg"):
            return self._extract_email_file(data, filename)

        # Images — return None so pipeline sends to VLM
        if mime_type.startswith("image/"):
            return None

        return None

    def _extract_text_file(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract text from plain text files."""
        try:
            text = data.decode("utf-8", errors="replace")
            if len(text) > 80000:
                text = text[:80000] + f"\n\n[... truncated, {len(text):,} total characters]"
            return ExtractedDocument(
                text=text,
                images=[],
                metadata={"filename": filename, "type": "text"},
            )
        except Exception:
            return ExtractedDocument(
                text="[Unable to decode file]",
                images=[],
                metadata={"filename": filename, "error": "decode_failed"},
            )

    def _extract_pdf(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract text and convert pages to images for VLM."""
        text_parts = []
        images = []
        page_count = 0
        is_scanned = False

        # Try text extraction first
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                page_count = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                    # Extract tables
                    for table in page.extract_tables():
                        table_text = "\n".join(
                            " | ".join(str(cell) if cell else "" for cell in row)
                            for row in table
                        )
                        text_parts.append(f"--- Table on Page {i+1} ---\n{table_text}")

                # If very little text extracted, likely scanned
                total_text = '\n'.join(text_parts)
                if len(total_text.strip()) < 100 and page_count > 0:
                    is_scanned = True

        except ImportError:
            text_parts.append("[pdfplumber not installed — cannot extract text]")
        except Exception as e:
            # Might be encrypted
            if "password" in str(e).lower() or "encrypted" in str(e).lower():
                return self._extract_encrypted_pdf(data, filename)
            text_parts.append(f"[PDF text extraction error: {e}]")

        # Convert pages to images for VLM (up to 10 pages)
        try:
            images = self._pdf_to_images(data, max_pages=10)
        except ImportError:
            pass  # pdf2image not installed
        except Exception:
            pass

        return ExtractedDocument(
            text='\n\n'.join(text_parts) if text_parts else "[No text extracted]",
            images=images,
            metadata={"filename": filename, "type": "pdf", "pages": page_count},
            is_scanned=is_scanned,
            page_count=page_count,
        )

    def _extract_encrypted_pdf(self, data: bytes, filename: str) -> ExtractedDocument:
        """Try to decrypt PDF with known passwords."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))

            if not reader.is_encrypted:
                return ExtractedDocument(
                    text="[PDF appears encrypted but pypdf couldn't process it]",
                    images=[],
                    metadata={"filename": filename, "error": "encryption_parse_error"},
                )

            # Try each password
            for password in self.pdf_passwords:
                try:
                    reader.decrypt(password)
                    # Success — extract text
                    text_parts = []
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"--- Page {i+1} ---\n{page_text}")

                    return ExtractedDocument(
                        text='\n\n'.join(text_parts) if text_parts else "[Decrypted but no text extracted]",
                        images=self._pdf_to_images(data, max_pages=10),
                        metadata={"filename": filename, "type": "pdf_encrypted", "decrypted": True},
                        page_count=len(reader.pages),
                    )
                except Exception:
                    continue

            # No password worked
            return ExtractedDocument(
                text=f"[Encrypted PDF — none of the {len(self.pdf_passwords)} stored passwords worked]",
                images=[],
                metadata={"filename": filename, "error": "encrypted_no_password"},
            )

        except ImportError:
            return ExtractedDocument(
                text="[Encrypted PDF — install pypdf for decryption]",
                images=[],
                metadata={"filename": filename, "error": "pypdf_not_installed"},
            )

    def _pdf_to_images(self, data: bytes, max_pages: int = 10) -> list[bytes]:
        """Convert PDF pages to PNG images for VLM processing."""
        from pdf2image import convert_from_bytes
        images = []
        pil_images = convert_from_bytes(data, dpi=150, first_page=1, last_page=max_pages)
        for img in pil_images:
            buf = io.BytesIO()
            img.save(buf, format="PNG", optimize=True)
            images.append(buf.getvalue())
        return images

    def _extract_docx(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract text from DOCX and convert to images for VLM."""
        text_parts = []
        images = []

        try:
            from docx import Document
            doc = Document(io.BytesIO(data))

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        text_parts.append(" | ".join(cells))

        except ImportError:
            text_parts.append("[python-docx not installed]")
        except Exception as e:
            text_parts.append(f"[DOCX extraction error: {e}]")

        # Convert to images (render via LibreOffice if available)
        try:
            images = self._docx_to_images(data)
        except Exception:
            pass

        return ExtractedDocument(
            text='\n'.join(text_parts) if text_parts else "[Empty document]",
            images=images,
            metadata={"filename": filename, "type": "docx"},
        )

    def _docx_to_images(self, data: bytes) -> list[bytes]:
        """Convert DOCX to images using LibreOffice."""
        import subprocess
        import tempfile

        with tempfile.TemporaryDirectory() as tmpdir:
            docx_path = os.path.join(tmpdir, "input.docx")
            with open(docx_path, "wb") as f:
                f.write(data)

            # Convert to PDF first, then to images
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", docx_path, "--outdir", tmpdir],
                capture_output=True, timeout=30,
            )

            pdf_path = os.path.join(tmpdir, "input.pdf")
            if os.path.exists(pdf_path):
                with open(pdf_path, "rb") as f:
                    return self._pdf_to_images(f.read(), max_pages=5)

        return []

    def _extract_xlsx(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract text from XLSX."""
        text_parts = []
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        text_parts.append(" | ".join(cells))
                        row_count += 1
                    if row_count >= 500:
                        text_parts.append(f"[... {ws.max_row - 500} more rows]")
                        break
            wb.close()
        except ImportError:
            text_parts.append("[openpyxl not installed]")
        except Exception as e:
            text_parts.append(f"[XLSX extraction error: {e}]")

        return ExtractedDocument(
            text='\n'.join(text_parts) if text_parts else "[Empty spreadsheet]",
            images=[],
            metadata={"filename": filename, "type": "xlsx"},
        )

    def _extract_pptx(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract text from PPTX."""
        text_parts = []
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
        except ImportError:
            text_parts.append("[python-pptx not installed]")
        except Exception as e:
            text_parts.append(f"[PPTX extraction error: {e}]")

        return ExtractedDocument(
            text='\n'.join(text_parts) if text_parts else "[Empty presentation]",
            images=[],
            metadata={"filename": filename, "type": "pptx"},
        )

    def _extract_zip(self, data: bytes, filename: str) -> ExtractedDocument:
        """List ZIP contents and extract small text files within."""
        text_parts = []
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                text_parts.append("ZIP archive contents:")
                for info in zf.infolist():
                    text_parts.append(f"  {info.filename} ({info.file_size:,} bytes)")

                    # If it's a small text file inside the zip, extract it
                    if info.file_size < 50000:
                        ext = os.path.splitext(info.filename.lower())[1]
                        if ext in self.TEXT_EXTENSIONS:
                            try:
                                content = zf.read(info.filename).decode("utf-8", errors="replace")
                                text_parts.append(f"  --- Content of {info.filename} ---")
                                text_parts.append(content[:5000])
                                text_parts.append(f"  --- End of {info.filename} ---")
                            except Exception:
                                pass
        except Exception as e:
            text_parts.append(f"[ZIP extraction error: {e}]")

        return ExtractedDocument(
            text='\n'.join(text_parts),
            images=[],
            metadata={"filename": filename, "type": "zip"},
        )

    def _extract_email_file(self, data: bytes, filename: str) -> ExtractedDocument:
        """Extract from .eml files."""
        text_parts = []
        try:
            from email import message_from_bytes
            msg = message_from_bytes(data)
            text_parts.append(f"Subject: {msg.get('subject', '')}")
            text_parts.append(f"From: {msg.get('from', '')}")
            text_parts.append(f"To: {msg.get('to', '')}")
            text_parts.append(f"Date: {msg.get('date', '')}")
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_parts.append(part.get_payload(decode=True).decode("utf-8", errors="replace"))
                        break
            else:
                text_parts.append(msg.get_payload(decode=True).decode("utf-8", errors="replace"))
        except Exception as e:
            text_parts.append(f"[EML extraction error: {e}]")

        return ExtractedDocument(
            text='\n'.join(text_parts),
            images=[],
            metadata={"filename": filename, "type": "eml"},
        )
